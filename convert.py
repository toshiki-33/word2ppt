"""
Word原稿→Claude.ai構成Markdown→PPTX 変換スクリプト

使い方:
    python convert.py input.md output.pptx
    python convert.py input.md output.pptx --template template/base_template.pptx

入力MDの想定フォーマット（Claude.ai Project が v4 インストラクションで出力するもの）:

    ---
    ## slide1: タイトル
    ### 本文
    本文中の穴埋めは「（数字　空白群　）」のまま白文字で表示

    ### 黄色字＝穴埋めの答え
    1. 氷期
    2. 間氷期
    ---

設計思想（v19・理想PPT「高２NO３（弥生）」リバースエンジニアリングに基づく）：
- 本文枠1つに全本文を白文字で配置（穴埋めは括弧+空白群のまま）
- 各黄色字を独立 textbox として絶対座標で配置
- 各黄色字シェイプにシェイプ単位の visibility set (Appear) クリックアニメ
- 位置は対応行のY座標 + 対応括弧のX座標で自動計算（妻が後で微調整する想定）
"""
import argparse
import math
import re
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt


FONT_NAME = "HGSゴシックM"
TITLE_FONT_PT = 60
# 理想PPT準拠: 本文36pt / 黄色字44pt
BODY_FONT_PT_DEFAULT = 36
YELLOW_FONT_PT_DEFAULT = 44

SLIDE_W_CM = 33.87
SLIDE_H_CM = 19.05

TITLE_BODY_X = 2.3
TITLE_BODY_Y = 3.6
TITLE_BODY_W = 29.2
TITLE_BODY_H = 14.3

# 本文枠（理想PPT準拠: pos=(0.37, 0.29) size=(33.50x16.63)）
BODY_X = 0.37
BODY_Y = 0.29
BODY_W = 33.50
BODY_H = 16.63

# 黄色字シェイプ（理想PPT準拠: 高さ2.14cm固定、幅は文字数依存）
YELLOW_SHAPE_HEIGHT_CM = 2.14
YELLOW_CHAR_WIDTH_CM = 1.85  # 44pt太字HGSゴシックM の実測寄り
YELLOW_PADDING_CM = 0.6  # 左右合計の余白

# 本文物理特性（36pt太字HGSゴシックM）
BODY_LINE_HEIGHT_CM = 1.10
BODY_CHAR_WIDTH_FULL_CM = 1.27  # 全角
BODY_CHAR_WIDTH_HALF_CM = 0.64  # 半角
BODY_MARGIN_TOP_CM = 0.05
BODY_MARGIN_LEFT_CM = 0.10


PT_NOTE_PATTERN = re.compile(r"\s*[（(]\s*\d+\s*pt\s*[）)]")
LIST_PREFIX_PATTERN = re.compile(r"^\s*[-・*]\s+")


def strip_noise(text):
    """構成MDに混入しがちなノイズを除去：pt表記・行頭リストハイフン"""
    text = PT_NOTE_PATTERN.sub("", text)
    text = LIST_PREFIX_PATTERN.sub("", text)
    return text


def parse_md(md_text):
    """構成MDをパースしてスライドのリストを返す（白文字本文＋黄色字のみ）"""
    slides = []
    blocks = re.split(r"^---\s*$", md_text, flags=re.MULTILINE)
    for block in blocks:
        block = block.strip()
        if not block:
            continue
        m = re.search(r"##\s*slide\s*(\d+)\s*[:：]?\s*(.*)", block, re.IGNORECASE)
        if not m:
            continue
        slide_no = int(m.group(1))
        slide_title = m.group(2).strip()

        body_lines = []
        yellow_items = []
        current = None
        in_meta_block = False

        for raw in block.split("\n"):
            line = raw.rstrip()
            if re.match(r"^\s*###\s*本文", line):
                current = "body"
                in_meta_block = False
                continue
            if re.match(r"^\s*###\s*(強調|赤字)", line):
                current = "ignore"
                in_meta_block = False
                continue
            if re.match(r"^\s*###\s*黄色字", line):
                current = "yellow"
                in_meta_block = False
                continue
            if re.match(r"^\s*##[^#]", line) or re.match(r"^\s*#[^#]", line):
                continue

            # メタブロック（【要確認】【追記提案】等）はスライドに含めない
            if re.match(r"^\s*【(要確認|追記提案|別説あり|スライド化困難|表記注意)】", line):
                in_meta_block = True
                continue
            if in_meta_block and not line.strip():
                in_meta_block = False
                continue
            if in_meta_block:
                continue

            if current == "body":
                stripped = line.strip()
                if stripped.startswith("[") and stripped.endswith("]"):
                    continue
                cleaned = strip_noise(line)
                # 空行はスキップ（妻の希望: フォントサイズ守りつつ詰め込み）
                if not cleaned.strip():
                    continue
                body_lines.append(cleaned)
            elif current == "yellow":
                cleaned = strip_noise(line)
                m2 = re.match(r"^\s*\d+\.\s*(.+)$", cleaned)
                if not m2:
                    # ハイフン記法の黄色字も受理（"- 寺院" 形式）
                    m2 = re.match(r"^\s*[-・*]\s*(.+)$", line)
                    if m2:
                        cleaned_item = strip_noise(m2.group(1))
                        cleaned_item = re.sub(r"^【要確認】\s*", "【要確認】", cleaned_item)
                        if cleaned_item.strip():
                            yellow_items.append(cleaned_item.strip())
                    continue
                txt = m2.group(1).strip()
                txt = re.sub(r"^【要確認】\s*", "【要確認】", txt)
                if txt:
                    yellow_items.append(txt)

        while body_lines and not body_lines[0].strip():
            body_lines.pop(0)
        while body_lines and not body_lines[-1].strip():
            body_lines.pop()

        slides.append(
            {
                "no": slide_no,
                "title": slide_title,
                "body": "\n".join(body_lines),
                "yellow": yellow_items,
            }
        )
    slides.sort(key=lambda s: s["no"])
    return slides


def remove_all_slides(prs):
    """テンプレの既存スライドを package ごと削除する"""
    sldIdLst = prs.slides._sldIdLst
    slide_id_elements = list(sldIdLst)
    for sldId in slide_id_elements:
        rId = sldId.get(qn("r:id"))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)

    pkg = prs.part.package
    for part in list(pkg.iter_parts()):
        if part.partname.startswith("/ppt/slides/slide"):
            try:
                pkg._parts.pop(part.partname, None)
            except Exception:
                pass


def set_black_background(slide):
    """スライド背景を黒にする（XML直接操作）"""
    spTree = slide.shapes._spTree
    cSld = spTree.getparent()
    for existing_bg in cSld.findall(qn("p:bg")):
        cSld.remove(existing_bg)

    bg_xml = (
        '<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        "<p:bgPr>"
        '<a:solidFill><a:srgbClr val="000000"/></a:solidFill>'
        '<a:effectLst/>'
        "</p:bgPr>"
        "</p:bg>"
    )
    bg_elem = etree.fromstring(bg_xml)
    cSld.insert(0, bg_elem)


def set_run_common(run, text, pt, color_rgb, bold=True):
    run.text = text
    run.font.size = Pt(pt)
    run.font.bold = bold
    run.font.name = FONT_NAME
    run.font.color.rgb = color_rgb


# 物理上限ベースの閾値（36pt太字HGSゴシックM、line_spacing=1.0）
# 本文枠 16.63cm / 行高 1.10cm ≒ 15行限界
# 1行幅 33.30cm / 全角 1.27cm ≒ 26字
CHARS_PER_LINE = 26
MAX_BODY_LINES_PER_SLIDE = 14
MAX_BODY_TOTAL_CHARS = 340
MAX_LINE_HARD_CHARS = 42
MAX_YELLOW_PER_SLIDE = 8


def count_wrapped_lines(text, chars_per_line=CHARS_PER_LINE):
    """折り返しと空行を考慮した実行行数"""
    total = 0
    for line in text.split("\n"):
        if line.strip():
            total += max(1, (len(line) + chars_per_line - 1) // chars_per_line)
        else:
            total += 1
    return total


def needs_split(body, yellow):
    """分割が必要か判定（行数・総字数・最長行・黄文字数の4軸）"""
    wrapped = count_wrapped_lines(body)
    if wrapped > MAX_BODY_LINES_PER_SLIDE:
        return True
    total_chars = sum(len(l) for l in body.split("\n"))
    if total_chars > MAX_BODY_TOTAL_CHARS:
        return True
    max_line = max((len(l) for l in body.split("\n")), default=0)
    if max_line > MAX_LINE_HARD_CHARS:
        return True
    if len(yellow) > MAX_YELLOW_PER_SLIDE:
        return True
    return False


def force_split_long_line(line, max_chars=MAX_LINE_HARD_CHARS):
    """1行が長すぎる場合、句点→読点→文字数の順で強制分割。改行なし長文の保険。"""
    if len(line) <= max_chars:
        return [line]
    pieces = [p for p in re.split(r"(?<=。)", line) if p]
    if len(pieces) > 1 and all(len(p) <= max_chars for p in pieces):
        return pieces
    result = []
    for p in pieces:
        if len(p) <= max_chars:
            result.append(p)
            continue
        sub = [s for s in re.split(r"(?<=[、，])", p) if s]
        if len(sub) > 1 and all(len(s) <= max_chars for s in sub):
            result.extend(sub)
            continue
        for s in sub or [p]:
            while len(s) > max_chars:
                result.append(s[:max_chars])
                s = s[max_chars:]
            if s:
                result.append(s)
    return result


def split_body_smart(body, max_lines=MAX_BODY_LINES_PER_SLIDE):
    """本文を意味の区切り優先で分割する。総字数も考慮し、長行は事前に強制分割。"""
    raw_lines = body.split("\n")
    lines = []
    for raw in raw_lines:
        if len(raw) > MAX_LINE_HARD_CHARS:
            lines.extend(force_split_long_line(raw))
        else:
            lines.append(raw)
    parts = []
    current = []

    def measure():
        joined = "\n".join(current)
        wrapped = count_wrapped_lines(joined)
        total_chars = sum(len(l) for l in joined.split("\n"))
        return wrapped, total_chars

    for line in lines:
        current.append(line)
        wrapped, total_chars = measure()
        over = wrapped > max_lines or total_chars > MAX_BODY_TOTAL_CHARS
        if over and len(current) > 1:
            cut = None
            for j in range(len(current) - 2, 0, -1):
                if not current[j].strip():
                    cut = j
                    break
                stripped = current[j].strip()
                if stripped.startswith(("◎", "※", "・")) or (
                    stripped and stripped[0] in "（("
                ):
                    cut = j
                    break
            if cut is None:
                cut = len(current) - 1

            head = current[:cut]
            while head and not head[-1].strip():
                head.pop()
            if head:
                parts.append("\n".join(head).strip())

            tail = current[cut:]
            while tail and not tail[0].strip():
                tail.pop(0)
            current = tail

    if current:
        leftover = "\n".join(current).strip()
        if leftover:
            parts.append(leftover)

    return parts or [body]


# 全括弧をいったん拾うパターン（中身判定は別関数で）
ANY_BRACKET_PATTERN = re.compile(r"[（(]([^（）()]*?)[）)]")
SPACE_RUN_PATTERN = re.compile(r"[\s　]{5,}")

# 行間の意味区切り判定用
NUMBER_HEAD_CHARS = "①②③④⑤⑥⑦⑧⑨⑩⑪⑫⑬⑭⑮"
SECTION_HEAD_CHARS = "◎○●・※→⇒"
SENTENCE_END_CHARS = "。！？"


def is_section_break(prev_line, curr_line):
    """直前の論理行と現在の論理行の間で行間を空けるべきか判定。
    意味の区切り（番号項目の開始・節記号・句点終わり）で True を返す。
    """
    if not prev_line or not curr_line:
        return False
    prev_stripped = prev_line.rstrip()
    curr_stripped = curr_line.lstrip()
    if not prev_stripped or not curr_stripped:
        return False
    # 現在の行頭が番号付き or 節記号 → 区切り
    if curr_stripped[0] in NUMBER_HEAD_CHARS or curr_stripped[0] in SECTION_HEAD_CHARS:
        return True
    # 現在の行頭が「数字.」or「数字)」 → 区切り
    if (
        curr_stripped[0].isdigit()
        and len(curr_stripped) > 1
        and curr_stripped[1] in ".．)）"
    ):
        return True
    # 直前が句点終わり → 区切り
    if prev_stripped[-1] in SENTENCE_END_CHARS:
        return True
    return False


def is_hole_bracket(content):
    """括弧の中身が穴埋めっぽいか判定。
    妻の原稿スタイル: 「（数字　空白群　）」or「（　空白群　）」or「（　空白群　＋固定テキスト）」
    """
    # 中身に空白を2個以上含む → 穴埋めスペース
    if re.search(r"[\s　]{2,}", content):
        return True
    # 数字+空白で始まる（番号付き穴埋め）
    if re.match(r"^\s*\d{1,2}[\s　]+", content):
        return True
    return False


def count_holes_in_text(text):
    """本文中の穴埋め個数（is_hole_bracket と同じ判定）。
    括弧で囲まれた穴埋め＋括弧なしの空白群5個以上もカウント。
    """
    n = 0
    for m in ANY_BRACKET_PATTERN.finditer(text):
        if is_hole_bracket(m.group(1)):
            n += 1
    text_no_brackets = ANY_BRACKET_PATTERN.sub("", text)
    n += len(re.findall(r"[\s　]{5,}", text_no_brackets))
    return n


def assign_yellow_to_parts(body_parts, yellow_items):
    """黄色字を本文パートの穴埋め個数に応じて出現順に割り振る"""
    result = [[] for _ in body_parts]
    y_idx = 0
    for p_idx, part in enumerate(body_parts):
        hole_count = count_holes_in_text(part)
        for _ in range(hole_count):
            if y_idx < len(yellow_items):
                result[p_idx].append(yellow_items[y_idx])
                y_idx += 1
    while y_idx < len(yellow_items):
        result[-1].append(yellow_items[y_idx])
        y_idx += 1
    return result


def split_yellow_if_overflow(slides_expanded):
    """黄色字が上限超のスライドを「本文も穴位置で対応分割」した上でチャンク分割。
    答えと本文の対応がズレないよう、本文を穴数MAX_YELLOW_PER_SLIDEごとに切る。
    """
    result = []
    for sd in slides_expanded:
        y = sd["yellow"]
        if len(y) <= MAX_YELLOW_PER_SLIDE:
            result.append(sd)
            continue
        body_lines = sd["body"].split("\n")
        body_chunks = []
        yellow_chunks = []
        current_lines = []
        current_yellow = []
        hole_in_chunk = 0
        y_idx = 0
        for line in body_lines:
            line_holes = count_holes_in_text(line)
            # この行を加えると上限超える かつ 既に何か入ってる → ここで切る
            if hole_in_chunk + line_holes > MAX_YELLOW_PER_SLIDE and current_lines:
                body_chunks.append("\n".join(current_lines))
                yellow_chunks.append(current_yellow)
                current_lines = []
                current_yellow = []
                hole_in_chunk = 0
            current_lines.append(line)
            for _ in range(line_holes):
                if y_idx < len(y):
                    current_yellow.append(y[y_idx])
                    y_idx += 1
            hole_in_chunk += line_holes
        if current_lines:
            body_chunks.append("\n".join(current_lines))
            yellow_chunks.append(current_yellow)
        # 残った黄色字（穴と対応取れなかった分）は最後のチャンクに追加
        while y_idx < len(y):
            yellow_chunks[-1].append(y[y_idx])
            y_idx += 1

        for idx, (bc, yc) in enumerate(zip(body_chunks, yellow_chunks)):
            title = sd["title"]
            if idx > 0 and not title.endswith("（続き）"):
                title = title + "（続き）"
            result.append(dict(sd, body=bc, yellow=yc, title=title))
    return result


def split_slide_if_needed(slide_data):
    """本文の物理量または黄色字個数が上限を超えたらスライド分割"""
    body = slide_data["body"]
    yellow = slide_data["yellow"]

    if not needs_split(body, yellow):
        return [slide_data]

    body_parts = split_body_smart(body) if body.strip() else [body]
    yellow_per_part = assign_yellow_to_parts(body_parts, yellow)

    expanded = []
    for idx, (bp, yp) in enumerate(zip(body_parts, yellow_per_part)):
        suffix = "" if idx == 0 else "（続き）"
        expanded.append(
            {
                "no": slide_data["no"],
                "title": slide_data["title"] + suffix,
                "body": bp,
                "yellow": yp,
            }
        )

    return split_yellow_if_overflow(expanded)


def is_title_slide(slide):
    if slide["no"] == 1:
        return True
    body = slide["body"].strip()
    if not slide["yellow"] and len(body) < 80 and "\n" not in body:
        return True
    return False


def add_title_slide(prs, slide_data):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_black_background(slide)
    tb = slide.shapes.add_textbox(
        Cm(TITLE_BODY_X), Cm(TITLE_BODY_Y), Cm(TITLE_BODY_W), Cm(TITLE_BODY_H)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    text = slide_data["body"].strip() or slide_data["title"]
    run = p.add_run()
    set_run_common(run, text, TITLE_FONT_PT, RGBColor(0xFF, 0xFF, 0xFF))
    return slide


def build_inline_segments(line, yellow_iter):
    """1論理行を [(text, color), ...] のセグメント列に変換し、
    各黄色ランの行内文字オフセット (st, end) も返す（半開区間）。
    color は "white" or "yellow"。
    """
    segments = []
    yellow_offsets = []

    # 括弧と長空白群を出現順に処理
    events = []
    for m in ANY_BRACKET_PATTERN.finditer(line):
        if is_hole_bracket(m.group(1)):
            events.append(("bracket", m.start(), m.end(), m.group(1)))

    # 括弧外の長空白群を検出するため、括弧位置をマスク
    masked = list(line)
    for _, s, e, _ in events:
        for i in range(s, e):
            masked[i] = "_"
    masked_str = "".join(masked)
    for m in SPACE_RUN_PATTERN.finditer(masked_str):
        events.append(("space", m.start(), m.end(), ""))
    events.sort(key=lambda x: x[1])

    cursor = 0  # 行内文字位置
    pos = 0    # 入力line上の処理済み位置
    for kind, s, e, inner in events:
        # 括弧前のテキスト（白）
        if pos < s:
            text = line[pos:s]
            segments.append((text, "white"))
            cursor += len(text)

        try:
            answer = next(yellow_iter)
        except StopIteration:
            # 答えが足りない: 元の括弧/空白をそのまま残す
            text = line[s:e]
            segments.append((text, "white"))
            cursor += len(text)
            pos = e
            continue

        if kind == "bracket":
            # 「（数字＋空白群＋固定テキスト）」から固定テキストを抽出
            fixed = re.sub(r"^\s*\d{1,2}[\s　]+", "", inner)
            fixed = re.sub(r"^[\s　]+", "", fixed)
            fixed = fixed.rstrip(" 　\t")
            # 開き括弧
            segments.append(("（", "white"))
            cursor += 1
            # 答え（黄色）
            yellow_st = cursor
            segments.append((answer, "yellow"))
            cursor += len(answer)
            yellow_end = cursor
            yellow_offsets.append((yellow_st, yellow_end))
            # 固定テキスト
            if fixed:
                segments.append((fixed, "white"))
                cursor += len(fixed)
            # 閉じ括弧
            segments.append(("）", "white"))
            cursor += 1
        else:  # space run
            yellow_st = cursor
            segments.append((answer, "yellow"))
            cursor += len(answer)
            yellow_end = cursor
            yellow_offsets.append((yellow_st, yellow_end))

        pos = e

    if pos < len(line):
        text = line[pos:]
        segments.append((text, "white"))
        cursor += len(text)

    return segments, yellow_offsets


def _create_body_textbox(slide):
    tb = slide.shapes.add_textbox(Cm(BODY_X), Cm(BODY_Y), Cm(BODY_W), Cm(BODY_H))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    # 理想PPT準拠: anchor=top
    tf.margin_left = Cm(BODY_MARGIN_LEFT_CM)
    tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(BODY_MARGIN_TOP_CM)
    tf.margin_bottom = Cm(0.05)
    return tb, tf


def char_width_cm(c):
    """1文字の物理幅(cm) - 36pt太字HGSゴシックM"""
    code = ord(c)
    if code < 0x80:
        return BODY_CHAR_WIDTH_HALF_CM
    if 0xFF61 <= code <= 0xFF9F:
        return BODY_CHAR_WIDTH_HALF_CM
    return BODY_CHAR_WIDTH_FULL_CM


def line_visual_width_cm(line):
    return sum(char_width_cm(c) for c in line)


def physical_rows_for(line):
    """1論理行が物理何行に折返されるか"""
    if not line.strip():
        return 1
    inner_w = BODY_W - 2 * BODY_MARGIN_LEFT_CM
    return max(1, math.ceil(line_visual_width_cm(line) / inner_w))


def physical_row_at_col(line, col):
    """行内のcol位置が物理何行目か（0始まり）"""
    inner_w = BODY_W - 2 * BODY_MARGIN_LEFT_CM
    width_so_far = sum(char_width_cm(c) for c in line[:col])
    return int(width_so_far / inner_w)


def x_in_physical_row_cm(line, col):
    """行内のcol位置が物理行内の何cm目か（折返し済み）"""
    inner_w = BODY_W - 2 * BODY_MARGIN_LEFT_CM
    width_so_far = sum(char_width_cm(c) for c in line[:col])
    return width_so_far - int(width_so_far / inner_w) * inner_w


def find_hole_columns(line):
    """1論理行内の穴埋め位置（文字オフセット）を出現順で返す → [col_offset, ...]
    括弧付き穴埋めは「（」の右隣（穴の中）、長空白群はその先頭位置を返す。
    """
    events = []
    for m in ANY_BRACKET_PATTERN.finditer(line):
        if is_hole_bracket(m.group(1)):
            # 「（」の右隣を黄色字配置位置とする
            events.append((m.start() + 1, m.start(), m.end(), "bracket"))

    # 括弧外の長空白群
    masked = list(line)
    for _, s, e, _ in events:
        for i in range(s, e):
            masked[i] = "_"
    masked_str = "".join(masked)
    for m in SPACE_RUN_PATTERN.finditer(masked_str):
        events.append((m.start(), m.start(), m.end(), "space"))

    events.sort(key=lambda x: x[1])
    return [col for col, _, _, _ in events]


def add_yellow_shape(slide, text, x_cm, y_cm):
    """黄色字を独立 textbox として絶対座標で配置。
    理想PPT準拠: 高さ2.14cm固定、幅は文字数×1.85cm+余白0.6cm。
    """
    text_clean = text.replace("【要確認】", "")
    char_count = len(text_clean)
    # 全角・半角混在を考慮した幅推定
    width = sum(
        YELLOW_CHAR_WIDTH_CM if ord(c) >= 0x80 and not (0xFF61 <= ord(c) <= 0xFF9F) else YELLOW_CHAR_WIDTH_CM * 0.55
        for c in text_clean
    ) + YELLOW_PADDING_CM
    # 画面右端を超えないようクランプ
    max_w = SLIDE_W_CM - x_cm - 0.3
    width = max(2.5, min(width, max_w))
    # X座標が右端を超えないようクランプ
    if x_cm + width > SLIDE_W_CM - 0.3:
        x_cm = max(0.3, SLIDE_W_CM - width - 0.3)
    # Y座標が下端を超えないようクランプ
    if y_cm + YELLOW_SHAPE_HEIGHT_CM > SLIDE_H_CM - 0.3:
        y_cm = SLIDE_H_CM - YELLOW_SHAPE_HEIGHT_CM - 0.3

    tb = slide.shapes.add_textbox(
        Cm(x_cm), Cm(y_cm), Cm(width), Cm(YELLOW_SHAPE_HEIGHT_CM)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Cm(0.05)
    tf.margin_right = Cm(0.05)
    tf.margin_top = Cm(0)
    tf.margin_bottom = Cm(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    set_run_common(run, text, YELLOW_FONT_PT_DEFAULT, RGBColor(0xFF, 0xFF, 0x00))
    return tb.shape_id


def add_body_and_yellow_overlays(slide, body_text, yellow_items):
    """本文枠（白文字、穴埋めは括弧+空白群のまま）+ 各黄色字を独立シェイプとして
    対応行のY座標・対応列のX座標に配置する。理想PPT「高２NO３（弥生）」準拠。
    返り値: yellow_shape_ids（アニメ順 = 出現順）
    """
    lines = body_text.split("\n") if body_text else [""]

    # 1. 本文枠を配置（穴埋め記号はそのまま白文字で表示）
    tb_a, tf_a = _create_body_textbox(slide)
    first = True
    prev_p = None
    for i, line in enumerate(lines):
        if first:
            p = tf_a.paragraphs[0]
            first = False
        else:
            p = tf_a.add_paragraph()
        p.line_spacing = 1.0
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        if prev_p is not None and is_section_break(lines[i - 1], line):
            prev_p.space_after = Pt(8)
        prev_p = p
        if not line.strip():
            continue
        run = p.add_run()
        set_run_common(run, line, BODY_FONT_PT_DEFAULT, RGBColor(0xFF, 0xFF, 0xFF))

    # 2. 各論理行の物理Y座標を計算しつつ、穴埋め位置に黄色字を配置
    yellow_shape_ids = []
    yellow_idx = 0
    n_yellows = len(yellow_items)
    cumulative_y = BODY_Y + BODY_MARGIN_TOP_CM
    SECTION_BREAK_CM = 0.28  # Pt(8) ≒ 0.28cm

    for i, line in enumerate(lines):
        line_top_y = cumulative_y
        if not line.strip():
            cumulative_y += BODY_LINE_HEIGHT_CM
            continue
        hole_columns = find_hole_columns(line)
        for col_offset in hole_columns:
            if yellow_idx >= n_yellows:
                break
            ans = yellow_items[yellow_idx]
            yellow_idx += 1
            # 折返し物理行を考慮: 穴がこの行の物理何行目か
            row_offset = physical_row_at_col(line, col_offset)
            x_offset_in_row_cm = x_in_physical_row_cm(line, col_offset)
            x_cm = BODY_X + BODY_MARGIN_LEFT_CM + x_offset_in_row_cm
            # 黄色字シェイプは対応物理行の少し上に被せる（44pt > 36pt なので上ずらし）
            y_cm = line_top_y + row_offset * BODY_LINE_HEIGHT_CM - 0.4
            if y_cm < 0.1:
                y_cm = 0.1
            sp_id = add_yellow_shape(slide, ans, x_cm, y_cm)
            yellow_shape_ids.append(sp_id)
        # 次行に進む（折返し物理行数を考慮）
        cumulative_y += BODY_LINE_HEIGHT_CM * physical_rows_for(line)
        if i + 1 < len(lines) and is_section_break(line, lines[i + 1]):
            cumulative_y += SECTION_BREAK_CM
        if yellow_idx >= n_yellows:
            break

    # 3. 余りの黄色字（穴埋めと対応取れなかった分）は右サイドに縦並びでフォールバック
    fallback_y = BODY_Y + 0.5
    while yellow_idx < n_yellows:
        ans = yellow_items[yellow_idx]
        yellow_idx += 1
        x_cm = SLIDE_W_CM - 8.0
        sp_id = add_yellow_shape(slide, ans, x_cm, fallback_y)
        yellow_shape_ids.append(sp_id)
        fallback_y += YELLOW_SHAPE_HEIGHT_CM + 0.2

    return yellow_shape_ids


def build_overlay_appear_timing_xml(shape_ids):
    """オーバーレイ枠（シェイプ単位）にクリックで順次 Appear するアニメ XML を構築。
    NO1_reference.pptx slide10 の構造を踏襲（presetID=1 / presetClass=entr / 文字単位ではなくシェイプ単位）。
    """
    if not shape_ids:
        return None

    par_blocks = []
    next_id = 3
    for spid in shape_ids:
        outer_id = next_id
        mid_id = next_id + 1
        click_id = next_id + 2
        set_id = next_id + 3
        next_id += 4

        par_blocks.append(
            f"""
                <p:par>
                  <p:cTn id="{outer_id}" fill="hold">
                    <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
                    <p:childTnLst>
                      <p:par>
                        <p:cTn id="{mid_id}" fill="hold">
                          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                          <p:childTnLst>
                            <p:par>
                              <p:cTn id="{click_id}" presetID="1" presetClass="entr" presetSubtype="0" fill="hold" grpId="0" nodeType="clickEffect">
                                <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                                <p:childTnLst>
                                  <p:set>
                                    <p:cBhvr>
                                      <p:cTn id="{set_id}" dur="1" fill="hold">
                                        <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                                      </p:cTn>
                                      <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
                                      <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                                    </p:cBhvr>
                                    <p:to><p:strVal val="visible"/></p:to>
                                  </p:set>
                                </p:childTnLst>
                              </p:cTn>
                            </p:par>
                          </p:childTnLst>
                        </p:cTn>
                      </p:par>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>"""
        )

    bld_entries = "\n    ".join(
        f'<p:bldP spid="{spid}" grpId="0"/>' for spid in shape_ids
    )

    timing_xml = f"""<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>{''.join(par_blocks)}
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst>
              <p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:prevCondLst>
            <p:nextCondLst>
              <p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:nextCondLst>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
  <p:bldLst>
    {bld_entries}
  </p:bldLst>
</p:timing>"""
    return timing_xml


def attach_yellow_animations(slide, shape_ids):
    """スライドに黄色字シェイプのクリック順次appearアニメを注入"""
    timing_xml = build_overlay_appear_timing_xml(shape_ids)
    if timing_xml is None:
        return
    sld = slide._element
    for old in sld.findall(qn("p:timing")):
        sld.remove(old)
    sld.append(etree.fromstring(timing_xml))


def add_content_slide(prs, slide_data):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_black_background(slide)
    yellow_ids = add_body_and_yellow_overlays(
        slide, slide_data["body"], slide_data["yellow"]
    )
    attach_yellow_animations(slide, yellow_ids)
    return slide


def convert_md_to_pptx(md_text, output, template_path=None):
    """構成MD文字列をPPTXに変換する。

    Args:
        md_text (str): 構成Markdown本文
        output (str | BinaryIO): 保存先パス or 書き込み可能なバイナリストリーム
        template_path (str, optional): テンプレpptxパス
    Returns:
        dict: {"original": 元スライド数, "expanded": 分割後スライド数}
    """
    tpl_path = template_path or str(
        Path(__file__).parent / "template" / "base_template.pptx"
    )
    prs = Presentation(tpl_path)
    remove_all_slides(prs)

    slides = parse_md(md_text)
    if not slides:
        raise ValueError("構成MDからスライドを検出できませんでした。フォーマットを確認してください。")

    expanded = []
    for sd in slides:
        if is_title_slide(sd):
            expanded.append(sd)
        else:
            expanded.extend(split_slide_if_needed(sd))

    for sd in expanded:
        if is_title_slide(sd):
            add_title_slide(prs, sd)
        else:
            add_content_slide(prs, sd)

    prs.save(output)
    return {"original": len(slides), "expanded": len(expanded)}


def main():
    ap = argparse.ArgumentParser(description="構成MD→PPTX 変換")
    ap.add_argument("input_md", help="構成Markdownファイル")
    ap.add_argument("output_pptx", help="出力PPTXパス")
    ap.add_argument(
        "--template",
        default=None,
        help="テンプレPPTX（省略時は template/base_template.pptx）",
    )
    args = ap.parse_args()

    md = Path(args.input_md).read_text(encoding="utf-8")
    result = convert_md_to_pptx(md, args.output_pptx, args.template)
    split_added = result["expanded"] - result["original"]
    print(
        f"saved: {args.output_pptx} ({result['expanded']} slides"
        + (f", {split_added}枚は自動分割" if split_added else "")
        + ")"
    )


if __name__ == "__main__":
    main()
